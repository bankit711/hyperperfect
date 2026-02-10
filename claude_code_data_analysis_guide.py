"""Generate a PowerPoint guide for using Claude Code for data analysis."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
BG_DARK = RGBColor(0x0F, 0x0F, 0x14)
BG_CARD = RGBColor(0x1A, 0x1A, 0x24)
ACCENT = RGBColor(0xD4, 0x7B, 0x2A)  # warm orange
ACCENT_LIGHT = RGBColor(0xE8, 0x9B, 0x4A)
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_GRAY = RGBColor(0xB0, 0xB0, 0xB8)
TEXT_LIGHT = RGBColor(0xE0, 0xE0, 0xE4)
CODE_BG = RGBColor(0x12, 0x12, 0x1A)
GREEN = RGBColor(0x4E, 0xC9, 0x8B)
BLUE = RGBColor(0x5B, 0x8D, 0xEF)
PURPLE = RGBColor(0x9B, 0x6E, 0xE7)


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, corner_radius=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if corner_radius is not None:
        shape.adjustments[0] = corner_radius
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=TEXT_LIGHT, bullet_color=ACCENT, spacing=Pt(8)):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
        # Bullet character
        pPr = p._pPr
        if pPr is None:
            from pptx.oxml.ns import qn
            pPr = p._p.get_or_add_pPr()
        from pptx.oxml.ns import qn
        buChar = pPr.makeelement(qn('a:buChar'), {'char': '▸'})
        # Remove existing bullets
        for existing in pPr.findall(qn('a:buChar')):
            pPr.remove(existing)
        for existing in pPr.findall(qn('a:buNone')):
            pPr.remove(existing)
        pPr.append(buChar)
        # Bullet color
        buClr = pPr.makeelement(qn('a:buClr'), {})
        srgb = buClr.makeelement(qn('a:srgbClr'), {'val': f'{bullet_color}'})
        buClr.append(srgb)
        for existing in pPr.findall(qn('a:buClr')):
            pPr.remove(existing)
        pPr.append(buClr)

    return txBox


def add_code_block(slide, left, top, width, height, code_text, font_size=12):
    """Add a styled code block."""
    shape = add_shape(slide, left, top, width, height, CODE_BG, corner_radius=0.03)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_right = Inches(0.3)
    tf.margin_top = Inches(0.2)
    tf.margin_bottom = Inches(0.2)

    lines = code_text.strip().split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.name = "Consolas"
        p.font.color.rgb = TEXT_LIGHT
        p.space_after = Pt(2)
    return shape


def add_accent_bar(slide, left, top, width=Inches(0.06), height=Inches(0.6), color=ACCENT):
    """Add a vertical accent bar."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_number_badge(slide, left, top, number, color=ACCENT):
    """Add a numbered circle badge."""
    size = Inches(0.5)
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_WHITE
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(2)


# ── Build Presentation ──────────────────────────────────────────────

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height
MARGIN = Inches(0.8)
CONTENT_W = SLIDE_W - 2 * MARGIN


# ── SLIDE 1: Title ──────────────────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, BG_DARK)

# Decorative top bar
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.08))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

# Title
add_text_box(slide, MARGIN, Inches(2.0), CONTENT_W, Inches(1.2),
             "Claude Code for Data Analysis", font_size=48, bold=True, color=TEXT_WHITE,
             alignment=PP_ALIGN.CENTER)

# Subtitle
add_text_box(slide, MARGIN, Inches(3.2), CONTENT_W, Inches(0.8),
             "A Practical Guide to AI-Powered Data Workflows", font_size=24,
             color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# Divider
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(5.5), Inches(4.3), Inches(2.3), Inches(0.04))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

# Tagline
add_text_box(slide, MARGIN, Inches(4.8), CONTENT_W, Inches(0.6),
             "From raw data to insights — using your terminal", font_size=18,
             color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)


# ── SLIDE 2: What Is Claude Code? ───────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, MARGIN, Inches(0.5), CONTENT_W, Inches(0.8),
             "What Is Claude Code?", font_size=36, bold=True, color=TEXT_WHITE)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                MARGIN, Inches(1.2), Inches(2), Inches(0.04))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

add_text_box(slide, MARGIN, Inches(1.6), Inches(5.5), Inches(1.0),
             "Claude Code is Anthropic's agentic CLI tool that lives in your terminal. "
             "It can read, write, and execute code — making it ideal for data analysis tasks "
             "that involve scripting, file manipulation, and iterative exploration.",
             font_size=17, color=TEXT_LIGHT)

# Feature cards
features = [
    ("Read & Write Files", "Directly reads CSVs, JSON,\nParquet, databases, and more", GREEN),
    ("Execute Code", "Runs Python, R, SQL, and\nshell commands in your env", BLUE),
    ("Iterate & Refine", "Ask follow-up questions,\nrefine analysis on the fly", PURPLE),
    ("Generate Visuals", "Creates charts, tables, and\nexport-ready reports", ACCENT),
]

card_w = Inches(2.7)
card_h = Inches(2.2)
gap = Inches(0.35)
start_x = MARGIN
start_y = Inches(3.2)

for i, (title, desc, color) in enumerate(features):
    x = start_x + i * (card_w + gap)
    card = add_shape(slide, x, start_y, card_w, card_h, BG_CARD, corner_radius=0.04)

    # Color top bar on card
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, start_y, card_w, Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    add_text_box(slide, x + Inches(0.25), start_y + Inches(0.4), card_w - Inches(0.5), Inches(0.5),
                 title, font_size=18, bold=True, color=color)
    add_text_box(slide, x + Inches(0.25), start_y + Inches(1.0), card_w - Inches(0.5), Inches(1.0),
                 desc, font_size=14, color=TEXT_GRAY)


# ── SLIDE 3: Getting Started ────────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, MARGIN, Inches(0.5), CONTENT_W, Inches(0.8),
             "Getting Started", font_size=36, bold=True, color=TEXT_WHITE)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                MARGIN, Inches(1.2), Inches(2), Inches(0.04))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

# Step 1
add_number_badge(slide, MARGIN, Inches(1.7), "1")
add_text_box(slide, MARGIN + Inches(0.7), Inches(1.7), Inches(4), Inches(0.4),
             "Install Claude Code", font_size=20, bold=True, color=TEXT_WHITE)
add_code_block(slide, MARGIN + Inches(0.7), Inches(2.2), Inches(5), Inches(0.6),
               "npm install -g @anthropic-ai/claude-code", font_size=14)

# Step 2
add_number_badge(slide, MARGIN, Inches(3.2), "2")
add_text_box(slide, MARGIN + Inches(0.7), Inches(3.2), Inches(4), Inches(0.4),
             "Navigate to your data project", font_size=20, bold=True, color=TEXT_WHITE)
add_code_block(slide, MARGIN + Inches(0.7), Inches(3.7), Inches(5), Inches(0.6),
               "cd ~/projects/my-analysis", font_size=14)

# Step 3
add_number_badge(slide, MARGIN, Inches(4.7), "3")
add_text_box(slide, MARGIN + Inches(0.7), Inches(4.7), Inches(4), Inches(0.4),
             "Launch Claude Code", font_size=20, bold=True, color=TEXT_WHITE)
add_code_block(slide, MARGIN + Inches(0.7), Inches(5.2), Inches(5), Inches(0.6),
               "claude", font_size=14)

# Tip box on the right
tip_box = add_shape(slide, Inches(7.5), Inches(1.7), Inches(4.8), Inches(4.5), BG_CARD, corner_radius=0.04)
add_accent_bar(slide, Inches(7.5), Inches(1.7), Inches(0.06), Inches(4.5), ACCENT)

add_text_box(slide, Inches(7.9), Inches(1.9), Inches(4), Inches(0.4),
             "💡  Pro Tips", font_size=20, bold=True, color=ACCENT)

tips = [
    "Launch from the directory containing your data files so Claude can see them",
    "Claude has access to your full local environment — Python, R, SQL tools, etc.",
    "Use a virtual environment to keep analysis dependencies isolated",
    "Add a CLAUDE.md file to give Claude context about your project and data schema",
]
add_bullet_list(slide, Inches(7.9), Inches(2.5), Inches(4.2), Inches(3.5),
                tips, font_size=14, color=TEXT_LIGHT)


# ── SLIDE 4: Core Workflow ──────────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, MARGIN, Inches(0.5), CONTENT_W, Inches(0.8),
             "The Data Analysis Workflow", font_size=36, bold=True, color=TEXT_WHITE)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                MARGIN, Inches(1.2), Inches(2), Inches(0.04))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

# Workflow steps as horizontal cards
steps = [
    ("1", "Load & Explore", "Ask Claude to read your\ndata and summarize its\nstructure, types, and\nbasic statistics.", GREEN),
    ("2", "Clean & Transform", "Handle missing values,\nfix types, filter rows,\nreshape data — all via\nnatural language.", BLUE),
    ("3", "Analyze", "Run statistical tests,\nbuild models, calculate\nmetrics, or write custom\nanalysis scripts.", PURPLE),
    ("4", "Visualize", "Generate matplotlib,\nseaborn, or plotly charts.\nExport as PNG, SVG, or\ninteractive HTML.", ACCENT),
    ("5", "Export & Report", "Save results to CSV,\nExcel, or generate a\nformatted report with\nfindings.", RGBColor(0xE0, 0x5E, 0x7E)),
]

card_w = Inches(2.15)
card_h = Inches(3.2)
gap = Inches(0.2)
start_x = MARGIN
start_y = Inches(1.8)

for i, (num, title, desc, color) in enumerate(steps):
    x = start_x + i * (card_w + gap)
    card = add_shape(slide, x, start_y, card_w, card_h, BG_CARD, corner_radius=0.04)

    # Number circle
    add_number_badge(slide, x + Inches(0.1), start_y + Inches(0.2), num, color)

    add_text_box(slide, x + Inches(0.1), start_y + Inches(0.85), card_w - Inches(0.2), Inches(0.4),
                 title, font_size=17, bold=True, color=color)
    add_text_box(slide, x + Inches(0.1), start_y + Inches(1.4), card_w - Inches(0.2), Inches(1.6),
                 desc, font_size=13, color=TEXT_GRAY)

    # Arrow between cards
    if i < len(steps) - 1:
        arrow_x = x + card_w + Inches(0.02)
        add_text_box(slide, arrow_x, start_y + Inches(1.2), Inches(0.2), Inches(0.4),
                     "→", font_size=20, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom note
add_text_box(slide, MARGIN, Inches(5.5), CONTENT_W, Inches(0.8),
             "Each step is conversational — ask follow-up questions, refine your approach, "
             "and iterate without leaving the terminal.",
             font_size=15, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)


# ── SLIDE 5: Example Prompts ────────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, MARGIN, Inches(0.5), CONTENT_W, Inches(0.8),
             "Example Prompts for Data Analysis", font_size=36, bold=True, color=TEXT_WHITE)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                MARGIN, Inches(1.2), Inches(2), Inches(0.04))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

# Two columns of prompt examples
left_prompts = [
    ("Exploration", [
        '"Read sales_data.csv and give me a summary of the columns, types, and any data quality issues"',
        '"How many unique customers do we have? What\'s the distribution of order values?"',
    ], GREEN),
    ("Cleaning", [
        '"Drop rows where revenue is negative, fill missing dates with the previous row\'s date"',
        '"Standardize the country column — map all variations to ISO country codes"',
    ], BLUE),
]

right_prompts = [
    ("Analysis", [
        '"Calculate month-over-month revenue growth and flag any months with >20% decline"',
        '"Run a cohort retention analysis grouped by signup quarter"',
    ], PURPLE),
    ("Visualization", [
        '"Create a line chart of monthly revenue with a 3-month moving average overlay"',
        '"Build a heatmap of feature correlations and save it as correlation_matrix.png"',
    ], ACCENT),
]

col_w = Inches(5.5)
y_pos = Inches(1.6)

for col_x, prompts in [(MARGIN, left_prompts), (Inches(7), right_prompts)]:
    y = y_pos
    for category, examples, color in prompts:
        add_accent_bar(slide, col_x, y, Inches(0.05), Inches(0.35), color)
        add_text_box(slide, col_x + Inches(0.2), y, Inches(3), Inches(0.35),
                     category, font_size=18, bold=True, color=color)
        y += Inches(0.5)
        for ex in examples:
            add_text_box(slide, col_x + Inches(0.2), y, col_w - Inches(0.3), Inches(0.65),
                         ex, font_size=13, color=TEXT_LIGHT, font_name="Consolas")
            y += Inches(0.75)
        y += Inches(0.25)


# ── SLIDE 6: Working with Python ────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, MARGIN, Inches(0.5), CONTENT_W, Inches(0.8),
             "Python + Claude Code", font_size=36, bold=True, color=TEXT_WHITE)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                MARGIN, Inches(1.2), Inches(2), Inches(0.04))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

add_text_box(slide, MARGIN, Inches(1.6), CONTENT_W, Inches(0.6),
             "Claude Code writes and executes Python scripts directly. It has full access "
             "to your installed packages — pandas, numpy, scikit-learn, matplotlib, and more.",
             font_size=16, color=TEXT_LIGHT)

# Left: example conversation
add_text_box(slide, MARGIN, Inches(2.5), Inches(5.5), Inches(0.4),
             "Example Interaction", font_size=18, bold=True, color=ACCENT)

convo = '''You: "Load the parquet file in /data and show me
     the top 10 customers by lifetime value"

Claude: Reads the file, writes a pandas script,
        executes it, and shows results:

┌────────────┬──────────┬───────────┐
│ customer   │ orders   │ LTV ($)   │
├────────────┼──────────┼───────────┤
│ Acme Corp  │ 147      │ 284,500   │
│ Globex     │ 203      │ 251,200   │
│ Initech    │ 98       │ 198,750   │
└────────────┴──────────┴───────────┘

+ saves the script as ltv_analysis.py'''

add_code_block(slide, MARGIN, Inches(3.0), Inches(5.8), Inches(3.8), convo, font_size=12)

# Right: key capabilities
add_text_box(slide, Inches(7.2), Inches(2.5), Inches(5), Inches(0.4),
             "What Claude Can Do", font_size=18, bold=True, color=ACCENT)

capabilities = [
    "Write pandas scripts for data manipulation and aggregation",
    "Use scikit-learn for ML models, clustering, and predictions",
    "Create publication-quality charts with matplotlib and seaborn",
    "Process large files efficiently with chunked reading",
    "Install missing packages automatically with pip",
    "Save reusable scripts for reproducible analysis",
    "Read SQL databases and run queries directly",
    "Handle multiple file formats: CSV, Parquet, Excel, JSON, HDF5",
]
add_bullet_list(slide, Inches(7.2), Inches(3.0), Inches(5.2), Inches(4.0),
                capabilities, font_size=14, color=TEXT_LIGHT, bullet_color=GREEN)


# ── SLIDE 7: Advanced Techniques ────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, MARGIN, Inches(0.5), CONTENT_W, Inches(0.8),
             "Advanced Techniques", font_size=36, bold=True, color=TEXT_WHITE)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                MARGIN, Inches(1.2), Inches(2), Inches(0.04))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

# Three cards
techniques = [
    ("CLAUDE.md for Context",
     "Add a CLAUDE.md file to your project root describing your data schema, "
     "business rules, and analysis conventions. Claude reads this automatically "
     "and applies your context to every interaction.",
     '''# CLAUDE.md
## Data Schema
- revenue.csv: daily revenue
  columns: date, product, amount
- customers.csv: customer dim table
## Conventions
- Fiscal year starts April 1
- All currency in USD''',
     GREEN),
    ("Multi-File Pipelines",
     "Ask Claude to build multi-step pipelines that chain together — "
     "cleaning, feature engineering, modeling, and reporting in one flow.",
     '''You: "Build a pipeline that:
1. Cleans the raw CSV
2. Engineers features
3. Trains a churn model
4. Exports predictions to Excel
5. Generates a summary report"

Claude: Creates 5 scripts + a
        run_pipeline.sh orchestrator''',
     BLUE),
    ("Iterative Exploration",
     "Use Claude's conversation memory to drill down into findings. "
     "Each follow-up question builds on previous context without re-loading data.",
     '''You: "What drove the Q3 revenue dip?"
Claude: [analyzes, shows breakdown]

You: "Drill into the Enterprise segment"
Claude: [filters, deeper analysis]

You: "Compare to Q3 last year"
Claude: [year-over-year comparison]

You: "Save this whole analysis as
      a Jupyter notebook"''',
     PURPLE),
]

card_w = Inches(3.7)
card_h = Inches(5.0)
gap = Inches(0.3)
start_x = MARGIN

for i, (title, desc, code, color) in enumerate(techniques):
    x = start_x + i * (card_w + gap)
    card = add_shape(slide, x, Inches(1.6), card_w, card_h, BG_CARD, corner_radius=0.04)
    add_accent_bar(slide, x, Inches(1.6), Inches(0.05), card_h, color)

    add_text_box(slide, x + Inches(0.25), Inches(1.8), card_w - Inches(0.4), Inches(0.4),
                 title, font_size=18, bold=True, color=color)
    add_text_box(slide, x + Inches(0.25), Inches(2.3), card_w - Inches(0.4), Inches(1.2),
                 desc, font_size=12, color=TEXT_GRAY)
    add_code_block(slide, x + Inches(0.15), Inches(3.7), card_w - Inches(0.3), Inches(2.7),
                   code, font_size=10)


# ── SLIDE 8: Tools & Extensions ─────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, MARGIN, Inches(0.5), CONTENT_W, Inches(0.8),
             "Extending Claude Code", font_size=36, bold=True, color=TEXT_WHITE)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                MARGIN, Inches(1.2), Inches(2), Inches(0.04))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

# Left column: MCP Servers
add_text_box(slide, MARGIN, Inches(1.6), Inches(5.5), Inches(0.4),
             "MCP Servers for Data", font_size=22, bold=True, color=GREEN)
add_text_box(slide, MARGIN, Inches(2.1), Inches(5.5), Inches(0.6),
             "Model Context Protocol servers extend Claude with specialized data capabilities:",
             font_size=15, color=TEXT_LIGHT)

mcp_items = [
    "Filesystem server — controlled access to data directories",
    "Database servers — connect to Postgres, SQLite, BigQuery directly",
    "Google Sheets — read and write spreadsheet data",
    "Custom servers — build connectors to your internal data systems",
]
add_bullet_list(slide, MARGIN, Inches(2.7), Inches(5.5), Inches(2.5),
                mcp_items, font_size=14, color=TEXT_LIGHT, bullet_color=GREEN)

add_code_block(slide, MARGIN, Inches(4.8), Inches(5.5), Inches(1.5),
               '# Add a Postgres MCP server\nclaude mcp add postgres npx -- -y @modelcontextprotocol/server-postgres\n\n# Add filesystem access\nclaude mcp add data-files npx -- -y @modelcontextprotocol/server-filesystem ~/data',
               font_size=11)

# Right column: Slash Commands & Hooks
add_text_box(slide, Inches(7.2), Inches(1.6), Inches(5), Inches(0.4),
             "Useful Features", font_size=22, bold=True, color=BLUE)

features_list = [
    "Extended thinking — Claude shows its reasoning for complex analysis",
    "Sub-agents — delegate parts of analysis to parallel workers",
    "Headless mode — run batch analyses from shell scripts",
    "Custom slash commands — create reusable analysis shortcuts",
    "Hooks — auto-run linting or tests after code generation",
    "Jupyter integration — work inside notebooks with Claude",
]
add_bullet_list(slide, Inches(7.2), Inches(2.2), Inches(5.2), Inches(3.0),
                features_list, font_size=14, color=TEXT_LIGHT, bullet_color=BLUE)

# Headless example
add_text_box(slide, Inches(7.2), Inches(4.8), Inches(5), Inches(0.4),
             "Batch Analysis Example", font_size=16, bold=True, color=BLUE)
add_code_block(slide, Inches(7.2), Inches(5.3), Inches(5.2), Inches(1.3),
               '# Run an analysis headlessly\nclaude -p "Analyze sales.csv: compute \\\n  monthly growth rates, flag anomalies, \\\n  save results to report.md" --allowedTools bash,write',
               font_size=11)


# ── SLIDE 9: Best Practices ─────────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, MARGIN, Inches(0.5), CONTENT_W, Inches(0.8),
             "Best Practices", font_size=36, bold=True, color=TEXT_WHITE)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                MARGIN, Inches(1.2), Inches(2), Inches(0.04))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

practices = [
    ("Be Specific with Prompts",
     "Instead of \"analyze this data,\" say \"calculate the 90-day rolling churn rate "
     "by customer segment and flag segments above 15%.\" Specificity yields better results.",
     "▶  Vague: \"Look at the data\"\n▶  Better: \"Show the distribution of order values by region, "
     "with median and P95 stats\"",
     GREEN),
    ("Verify Results",
     "Always spot-check Claude's outputs. Ask it to show intermediate steps, "
     "validate row counts, and cross-check totals against known values.",
     "▶  \"Show me the row count before and after each filter step\"\n"
     "▶  \"Cross-check the total revenue against the summary table\"",
     BLUE),
    ("Save & Version Scripts",
     "Ask Claude to save analysis scripts to files so you can re-run, "
     "share, and version control your work. Reproducibility matters.",
     "▶  \"Save this analysis as a Python script I can re-run\"\n"
     "▶  \"Add docstrings explaining each transformation step\"",
     PURPLE),
    ("Use CLAUDE.md Effectively",
     "Document your data schema, naming conventions, and business rules "
     "in a CLAUDE.md file. This eliminates repetitive explanations.",
     "▶  Include column descriptions, valid ranges, and known quirks\n"
     "▶  Specify preferred chart styles, date formats, and output locations",
     ACCENT),
]

card_w = Inches(5.7)
card_h = Inches(2.2)
gap_x = Inches(0.4)
gap_y = Inches(0.3)

for i, (title, desc, example, color) in enumerate(practices):
    col = i % 2
    row = i // 2
    x = MARGIN + col * (card_w + gap_x)
    y = Inches(1.6) + row * (card_h + gap_y)

    card = add_shape(slide, x, y, card_w, card_h, BG_CARD, corner_radius=0.04)
    add_accent_bar(slide, x, y, Inches(0.05), card_h, color)

    add_text_box(slide, x + Inches(0.25), y + Inches(0.15), card_w - Inches(0.4), Inches(0.35),
                 title, font_size=17, bold=True, color=color)
    add_text_box(slide, x + Inches(0.25), y + Inches(0.55), card_w - Inches(0.4), Inches(0.65),
                 desc, font_size=12, color=TEXT_GRAY)
    add_code_block(slide, x + Inches(0.15), y + Inches(1.3), card_w - Inches(0.3), Inches(0.8),
                   example, font_size=10)


# ── SLIDE 10: Comparison — Claude Code vs ChatGPT ───────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, MARGIN, Inches(0.5), CONTENT_W, Inches(0.8),
             "Claude Code vs ChatGPT for Data Analysis", font_size=36, bold=True, color=TEXT_WHITE)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                MARGIN, Inches(1.2), Inches(2), Inches(0.04))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

# Table header row
table_top = Inches(1.6)
col_label_w = Inches(3.2)
col_claude_w = Inches(4.5)
col_gpt_w = Inches(4.5)
row_h = Inches(0.55)

# Header backgrounds
hdr_y = table_top
add_shape(slide, MARGIN, hdr_y, col_label_w, row_h, BG_CARD, corner_radius=0.02)
claude_hdr = add_shape(slide, MARGIN + col_label_w + Inches(0.1), hdr_y,
                        col_claude_w, row_h, RGBColor(0x1A, 0x2A, 0x1A), corner_radius=0.02)
gpt_hdr = add_shape(slide, MARGIN + col_label_w + col_claude_w + Inches(0.2), hdr_y,
                      col_gpt_w, row_h, RGBColor(0x2A, 0x1A, 0x1A), corner_radius=0.02)

add_text_box(slide, MARGIN + Inches(0.2), hdr_y + Inches(0.08), col_label_w, Inches(0.4),
             "Capability", font_size=15, bold=True, color=TEXT_GRAY)
add_text_box(slide, MARGIN + col_label_w + Inches(0.3), hdr_y + Inches(0.08),
             col_claude_w, Inches(0.4),
             "Anthropic  ·  Claude Code", font_size=15, bold=True, color=GREEN)
add_text_box(slide, MARGIN + col_label_w + col_claude_w + Inches(0.4), hdr_y + Inches(0.08),
             col_gpt_w, Inches(0.4),
             "OpenAI  ·  ChatGPT", font_size=15, bold=True, color=RGBColor(0xE0, 0x5E, 0x7E))

# Comparison rows
rows = [
    ("Environment",
     "Runs in your local terminal with full access to your files, tools, and environment",
     "Browser-based sandbox with file upload; no access to local environment"),
    ("Code Execution",
     "Executes code locally — use your own Python env, packages, databases, and hardware",
     "Runs code in a remote sandbox with pre-installed packages; limited customization"),
    ("Data Access",
     "Reads files directly from disk — no upload limits, handles large datasets natively",
     "Requires file upload (capped at ~500 MB); data leaves your machine"),
    ("Privacy",
     "Data stays on your machine; code runs locally without sending file contents to the cloud",
     "Files are uploaded to OpenAI servers for processing"),
    ("Extensibility",
     "MCP servers, custom hooks, headless/batch mode, CI/CD integration, shell scripting",
     "GPTs and plugins for web-based workflows; limited programmatic integration"),
    ("Workflow",
     "Terminal-native: integrates with git, bash, cron, scripts, and existing dev workflows",
     "GUI-based: conversational web UI with visual chart rendering"),
    ("Models",
     "Claude Opus 4.5, Sonnet 4, Haiku — optimized for code generation and analysis",
     "GPT-4o, o1, o3 — general-purpose with code interpreter capabilities"),
    ("Visualization",
     "Generates chart files (PNG, SVG, HTML) using your installed libraries",
     "Renders interactive charts inline in the browser"),
    ("Pricing",
     "API usage-based pricing; included with Claude Pro / Team / Enterprise plans",
     "ChatGPT Plus ($20/mo), Team ($25/mo), or API usage-based pricing"),
]

for i, (label, claude_val, gpt_val) in enumerate(rows):
    y = table_top + (i + 1) * (row_h + Inches(0.06))
    bg_color = BG_CARD if i % 2 == 0 else RGBColor(0x16, 0x16, 0x20)

    add_shape(slide, MARGIN, y, col_label_w, row_h, bg_color, corner_radius=0.01)
    add_shape(slide, MARGIN + col_label_w + Inches(0.1), y, col_claude_w, row_h, bg_color, corner_radius=0.01)
    add_shape(slide, MARGIN + col_label_w + col_claude_w + Inches(0.2), y, col_gpt_w, row_h, bg_color, corner_radius=0.01)

    add_text_box(slide, MARGIN + Inches(0.2), y + Inches(0.08), col_label_w - Inches(0.3), Inches(0.4),
                 label, font_size=12, bold=True, color=ACCENT)
    add_text_box(slide, MARGIN + col_label_w + Inches(0.3), y + Inches(0.05),
                 col_claude_w - Inches(0.3), Inches(0.45),
                 claude_val, font_size=11, color=TEXT_LIGHT)
    add_text_box(slide, MARGIN + col_label_w + col_claude_w + Inches(0.4), y + Inches(0.05),
                 col_gpt_w - Inches(0.3), Inches(0.45),
                 gpt_val, font_size=11, color=TEXT_LIGHT)


# ── SLIDE 11: Closing ───────────────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

# Decorative bottom bar
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                0, SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

add_text_box(slide, MARGIN, Inches(2.2), CONTENT_W, Inches(1.0),
             "Start Analyzing", font_size=48, bold=True, color=TEXT_WHITE,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, MARGIN, Inches(3.3), CONTENT_W, Inches(0.6),
             "Install Claude Code and turn your terminal into a data analysis powerhouse.",
             font_size=20, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# Divider
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(5.5), Inches(4.2), Inches(2.3), Inches(0.04))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

add_code_block(slide, Inches(4), Inches(4.8), Inches(5.3), Inches(0.7),
               "npm install -g @anthropic-ai/claude-code && claude", font_size=16)

add_text_box(slide, MARGIN, Inches(6.0), CONTENT_W, Inches(0.5),
             "docs.anthropic.com/claude-code", font_size=16, color=TEXT_GRAY,
             alignment=PP_ALIGN.CENTER)


# ── Save ─────────────────────────────────────────────────────────────

output_path = "/Users/davidingraham/hyperperfect/Claude_Code_Data_Analysis_Guide.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
